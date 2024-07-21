import pdfplumber
import docx2txt
from pptx import Presentation

def extract_text_from_pdf(pdf_path):
    with pdfplumber.open(pdf_path) as pdf:
        text = ""
        for page in pdf.pages:
            text += page.extract_text() + "\n"
    
    return text

def extract_text_from_word_document(word_path):
    text = docx2txt.process(word_path)
    
    return text

def extract_text_from_ppt(ppt_path):
    prs = Presentation(ppt_path)
    text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    
    return text